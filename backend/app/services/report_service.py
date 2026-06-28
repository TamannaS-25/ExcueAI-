import os
import datetime
from typing import Dict, Any, List

# PPTX libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ReportLab PDF libraries
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

# DOCX libraries
import docx
from docx.shared import Inches as DocxInches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Ensure temporary exports directory exists
EXPORTS_DIR = os.path.abspath("./exports")
os.makedirs(EXPORTS_DIR, exist_ok=True)

class NumberedCanvas(canvas.Canvas):
    """
    Custom canvas to calculate and render running headers and footers with page counts
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_page_decorations(self, page_count):
        # Suppress header/footer on title page
        if self._pageNumber == 1:
            return
            
        self.saveState()
        self.setFont("Helvetica", 9)
        self.setFillColor(colors.HexColor("#4b5563"))
        
        # Header
        self.drawString(54, 750, "ExcueAI - Enterprise Report Workspace")
        self.setStrokeColor(colors.HexColor("#e5e7eb"))
        self.setLineWidth(0.5)
        self.line(54, 742, 558, 742)
        
        # Footer
        page_str = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 40, page_str)
        self.drawString(54, 40, "CONFIDENTIAL - Internal Use Only")
        self.line(54, 52, 558, 52)
        
        self.restoreState()


def generate_pptx_report(topic: str, project_name: str = "", meeting_summary: str = "", document_context: str = "") -> str:
    prs = Presentation()
    
    # Theme Color definitions (Royal Indigo theme)
    c_dark_navy = RGBColor(15, 23, 42)
    c_indigo = RGBColor(99, 102, 241)
    c_slate = RGBColor(71, 85, 105)
    c_light_gray = RGBColor(241, 245, 249)
    c_white = RGBColor(255, 255, 255)
    
    # Helper to format slide title
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Georgia'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = c_indigo
        
    # Content structure
    slides_data = [
        {"title": "Introduction", "body": f"Introducing our analysis for project topic: {topic}. This review details background objectives and frameworks."},
        {"title": "Objectives", "body": f"• Establish a systematic methodology for addressing: {topic}.\n• Define core metrics and deliverables for Project {project_name}.\n• Streamline information transfers across department boundaries."},
        {"title": "Problem Statement", "body": f"Organizations experience efficiency bottlenecks due to fragmented operational documents and data silos. For project: {topic}, these limitations prevent real-time strategy adjustments and cause resource leakage."},
        {"title": "Analysis", "body": f"Based on historical metrics and notes:\n• Task completion durations vary by department.\n• RAG vector search matches documents context files.\n• Data: {document_context[:200]}..."},
        {"title": "Findings", "body": f"• Meeting summary notes highlight: {meeting_summary[:200]}...\n• Employees report communication overlaps.\n• Systemic data availability delays resolution speeds by 15%."},
        {"title": "Recommendations", "body": "• Integrate automated AI workspace handlers to digest meeting notes.\n• Migrate document knowledge bases into secure vector database.\n• Implement strict role-based permission rules to protect organizational details."},
        {"title": "Conclusion", "body": f"By deploying ExcueAI workflows for {topic}, the enterprise can improve productivity compliance, minimize review timelines, and enable quick query answers."}
    ]
    
    # 1. Slide 1: Title Slide (Dark Background)
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Fill background
    bg = slide.shapes.add_shape(
        1, # Rectangle type
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_dark_navy
    bg.line.color.rgb = c_dark_navy
    
    # Add title text box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = topic.upper()
    p.alignment = PP_ALIGN.LEFT
    p.font.name = 'Georgia'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_white
    
    p2 = tf.add_paragraph()
    p2.text = f"Project Workspace Analysis | {project_name}"
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.color.rgb = c_indigo
    p2.space_before = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = f"Generated by ExcueAI on {datetime.date.today().strftime('%B %d, %Y')}"
    p3.font.name = 'Arial'
    p3.font.size = Pt(12)
    p3.font.color.rgb = c_slate
    p3.space_before = Pt(30)
    
    # 2. Content Slides
    for data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide, data["title"])
        
        # Body content box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        lines = data["body"].split("\n")
        for idx, line in enumerate(lines):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = line
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = c_dark_navy
            p.space_after = Pt(14)
            if line.strip().startswith("•"):
                p.level = 0
                
    filename = f"PPT_Report_{int(datetime.datetime.utcnow().timestamp())}.pptx"
    filepath = os.path.join(EXPORTS_DIR, filename)
    prs.save(filepath)
    return filepath


def generate_pdf_report(topic: str, project_name: str = "", meeting_summary: str = "", document_context: str = "") -> str:
    filename = f"PDF_Report_{int(datetime.datetime.utcnow().timestamp())}.pdf"
    filepath = os.path.join(EXPORTS_DIR, filename)
    
    # Setup document
    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        leftMargin=54,
        rightMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = getSampleStyleSheet()
    
    # Modify default styles or create custom ones
    primary_color = colors.HexColor("#1e3a8a") # Deep Blue
    secondary_color = colors.HexColor("#4f46e5") # Indigo
    body_color = colors.HexColor("#1f2937") # Charcoal
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=28,
        leading=34,
        textColor=primary_color,
        spaceAfter=15
    )
    
    subtitle_style = ParagraphStyle(
        'DocSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor("#4b5563"),
        spaceAfter=30
    )
    
    h1_style = ParagraphStyle(
        'DocH1',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=secondary_color,
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=10,
        leading=15,
        textColor=body_color,
        spaceAfter=12
    )

    meta_style = ParagraphStyle(
        'DocMeta',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=9,
        textColor=colors.HexColor("#6b7280")
    )
    
    story = []
    
    # 1. Cover Info / Header Table
    story.append(Spacer(1, 40))
    story.append(Paragraph(topic.upper(), title_style))
    story.append(Paragraph(f"Workspace Executive Assessment Report | Project {project_name}", subtitle_style))
    
    # Metadata Box Table
    meta_data = [
        [Paragraph("AUTHOR:", meta_style), Paragraph("ExcueAI Intelligence Engine", body_style)],
        [Paragraph("DATE:", meta_style), Paragraph(datetime.date.today().strftime("%B %d, %Y"), body_style)],
        [Paragraph("PROJECT ID:", meta_style), Paragraph(f"EXC-{project_name.upper() or 'N/A'}", body_style)],
        [Paragraph("CLASSIFICATION:", meta_style), Paragraph("CONFIDENTIAL / COMPANY INTERNAL", body_style)]
    ]
    meta_table = Table(meta_data, colWidths=[120, 384])
    meta_table.setStyle(TableStyle([
        ('LINEBELOW', (0,0), (-1,-1), 0.5, colors.HexColor("#e5e7eb")),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ]))
    story.append(meta_table)
    story.append(Spacer(1, 30))
    story.append(PageBreak()) # Push actual content to page 2 for clean cover format
    
    # 2. Executive Summary
    story.append(Paragraph("Executive Summary", h1_style))
    story.append(Paragraph(
        f"This document presents the strategic and technical assessment regarding: {topic}. "
        "Through data aggregation and task metrics analysis, we identify optimization areas "
        "to decrease cycle latencies and support productivity.",
        body_style
    ))
    
    # 3. Objectives
    story.append(Paragraph("Objectives", h1_style))
    story.append(Paragraph(
        f"The primary objectives for project '{project_name}' include the following dimensions:",
        body_style
    ))
    story.append(Paragraph("• Coordinate cross-department resource mappings to avoid code duplication.", body_style))
    story.append(Paragraph("• Enforce task execution milestones and track deadline compliance values.", body_style))
    story.append(Paragraph("• Provide robust RAG data sources to minimize information search delays.", body_style))
    
    # 4. Analysis
    story.append(Paragraph("Analysis", h1_style))
    story.append(Paragraph(
        f"Our database analyzer evaluated document parameters. "
        f"Referenced contextual text segments: {document_context[:500]}...",
        body_style
    ))
    
    # 5. Findings
    story.append(Paragraph("Findings", h1_style))
    story.append(Paragraph(
        f"A review of meeting discussions highlights the following takeaways: "
        f"{meeting_summary[:500]}...",
        body_style
    ))
    
    # 6. Recommendations & Action Items
    story.append(Paragraph("Recommendations", h1_style))
    story.append(Paragraph(
        "Based on these metrics, we recommend: "
        "1. Migrate standard manuals and policies to local vector indices. "
        "2. Automate task boards scheduling routines using priority filters. "
        "3. Configure JWT expiration values and verify security compliance indices.",
        body_style
    ))
    
    # 7. Conclusion
    story.append(Paragraph("Conclusion", h1_style))
    story.append(Paragraph(
        f"Deploying these automated policies for {topic} ensures robust task management "
        "and builds a scalable foundation for upcoming project iterations.",
        body_style
    ))
    
    # Build PDF
    doc.build(story, canvasmaker=NumberedCanvas)
    return filepath


def generate_docx_report(topic: str, project_name: str = "", meeting_summary: str = "", document_context: str = "") -> str:
    doc = docx.Document()
    
    # Set document margins
    for section in doc.sections:
        section.top_margin = DocxInches(1)
        section.bottom_margin = DocxInches(1)
        section.left_margin = DocxInches(1)
        section.right_margin = DocxInches(1)
        
    # Title
    title = doc.add_heading(level=0)
    run = title.add_run(topic.upper())
    run.font.name = 'Calibri'
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = docx.shared.RGBColor(30, 58, 138) # Deep Blue
    
    # Subtitle
    sub = doc.add_paragraph()
    r_sub = sub.add_run(f"Project Assessment Report: {project_name}\nGenerated by ExcueAI on {datetime.date.today().strftime('%B %d, %Y')}\nCONFIDENTIAL")
    r_sub.font.size = Pt(11)
    r_sub.italic = True
    r_sub.font.color.rgb = docx.shared.RGBColor(107, 114, 128)
    
    # Divider line
    doc.add_paragraph("__________________________________________________________________")
    
    def add_styled_heading(text, level):
        h = doc.add_heading(text, level=level)
        for r in h.runs:
            r.font.name = 'Calibri'
            r.font.bold = True
            r.font.color.rgb = docx.shared.RGBColor(79, 70, 229) # Indigo
        return h
        
    # Executive Summary
    add_styled_heading("Executive Summary", 1)
    p = doc.add_paragraph()
    p.add_run(f"This report outlines the structural framework and analytics reviews concerning: {topic}. "
              f"We analyzed active tasks, meeting briefs, and uploaded files to propose optimizations.")
              
    # Objectives
    add_styled_heading("Objectives", 1)
    p = doc.add_paragraph("Key priorities:")
    doc.add_paragraph("• Audit current file storage and indexing pipelines.", style='List Bullet')
    doc.add_paragraph("• Resolve department task assignment bottlenecks.", style='List Bullet')
    doc.add_paragraph("• Enhance access validation rules using role mappings.", style='List Bullet')

    # Analysis
    add_styled_heading("Analysis", 1)
    p = doc.add_paragraph()
    p.add_run("The engineering and operations datasets were parsed. Document knowledge highlights:\n")
    p.add_run(document_context)
    
    # Findings
    add_styled_heading("Findings", 1)
    p = doc.add_paragraph()
    p.add_run(f"The meeting intelligence analyzer parsed transcripts and extracted summaries: {meeting_summary}")
    
    # Recommendations
    add_styled_heading("Recommendations", 1)
    doc.add_paragraph("• Deploy in-memory vector stores for rapid document retrieval.", style='List Number')
    doc.add_paragraph("• Automate tasks deadlines tracking dashboards.", style='List Number')
    doc.add_paragraph("• Integrate comparison scripts for tracking drafts modifications.", style='List Number')
    
    # Conclusion
    add_styled_heading("Conclusion", 1)
    p = doc.add_paragraph()
    p.add_run(f"Implementing the recommended steps for {topic} ensures efficiency improvements and "
              "resolves communication bottlenecks between managers and employees.")
              
    filename = f"Word_Report_{int(datetime.datetime.utcnow().timestamp())}.docx"
    filepath = os.path.join(EXPORTS_DIR, filename)
    doc.save(filepath)
    return filepath
